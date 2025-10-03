import io
import logging
import re
import asyncio
from aiogram import Bot, Dispatcher, types
from aiogram.filters import Command
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt, Inches

API_TOKEN = "ТВОЙ_ТОКЕН"  # заменишь на Railway через os.getenv потом

logging.basicConfig(level=logging.INFO)

bot = Bot(token=API_TOKEN)
dp = Dispatcher()

# === Транслитерация армянского текста в латиницу ===
arm_to_lat = {
    "ա": "a", "բ": "b", "գ": "g", "դ": "d", "ե": "e", "զ": "z", "է": "e",
    "ը": "y", "թ": "t", "ժ": "zh", "ի": "i", "լ": "l", "խ": "kh", "ծ": "ts",
    "կ": "k", "հ": "h", "ձ": "dz", "ղ": "gh", "ճ": "ch", "մ": "m", "յ": "y",
    "ն": "n", "շ": "sh", "ո": "o", "չ": "ch", "պ": "p", "ջ": "j", "ռ": "r",
    "ս": "s", "վ": "v", "տ": "t", "ր": "r", "ց": "ts", "ու": "u", "փ": "p",
    "ք": "q", "և": "ev", "օ": "o", "ֆ": "f",
}

def transliterate(text: str) -> str:
    text = text.lower()
    for arm, lat in arm_to_lat.items():
        text = text.replace(arm, lat)
    text = re.sub(r"[^a-z0-9\s]", "", text)
    return text.strip()

# === Очистка строки: убираем только знаки препинания в конце ===
def clean_line(line: str) -> str:
    return line.rstrip(".,!?;:…").strip()

def is_armenian(text: str) -> bool:
    return bool(re.search(r"[Ա-Ֆա-ֆ]", text))

# === Динамический подбор размера шрифта ===
def get_dynamic_font_size(line: str, is_title=False) -> Pt:
    base_size = 120 if is_title else 100
    length = len(line)

    if length < 15:
        return Pt(base_size)
    elif length < 30:
        return Pt(base_size - 20)
    elif length < 50:
        return Pt(base_size - 40)
    else:
        return Pt(base_size - 60)

# === Создание презентации ===
def create_presentation(text: str) -> (io.BytesIO, str):
    prs = Presentation()
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    lines = [clean_line(line) for line in text.split("\n") if line.strip()]
    if not lines:
        lines = ["ПУСТОЙ ТЕКСТ"]

    # === Название файла ===
    first_line = lines[0]
    if is_armenian(first_line):
        filename = transliterate(first_line)
    else:
        filename = first_line
    filename = re.sub(r"[^\w\sА-Яа-яЁёA-Za-z0-9]", "", filename).strip()
    filename = filename + ".pptx"

    for i, line in enumerate(lines):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
        tf = txBox.text_frame
        tf.clear()
        p = tf.add_paragraph()

        if i == 0:  # заголовок
            p.text = line.upper()  # всегда большими буквами
            p.font.size = get_dynamic_font_size(line, is_title=True)
        else:
            p.text = line
            p.font.size = get_dynamic_font_size(line, is_title=False)

        p.font.name = "Sylfaen"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        tf.word_wrap = False  # не переносим строки, пусть будет в одну
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # всегда сверху

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output, filename

# === Хэндлеры ===
@dp.message(Command("start"))
async def start_command(message: types.Message):
    await message.answer("Привет 👋 Отправь мне текст.\nПервая строка = заголовок (в одну строку и имя файла).")

@dp.message()
async def create_pptx(message: types.Message):
    text = message.text
    pptx_file, filename = create_presentation(text)
    await message.answer_document(
        document=types.BufferedInputFile(pptx_file.read(), filename=filename)
    )

# === Запуск ===
async def main():
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
